"""
SmartArtAgent — Draws SmartArt-style diagrams directly on PowerPoint slides using python-pptx.
Supports: chevron process flows, pyramid diagrams, circular cycle diagrams, timeline roadmaps.
Uses native pptx shapes (no images) for crisp, editable vector output.
"""

import json
import logging
from pathlib import Path
from typing import Optional, List, Dict, Any
from .base_agent import BaseAgent, AgentResult

logger = logging.getLogger(__name__)

# Content type → diagram type mapping
CONTENT_TYPE_MAP = {
    "process": "chevron",
    "chevron_process": "chevron",
    "process_flow": "chevron",
    "process_diagram": "chevron",
    "hierarchy": "org",
    "org_chart": "org",
    "strategy": "pyramid",
    "pyramid": "pyramid",
    "funnel_diagram": "pyramid",
    "timeline": "roadmap",
    "roadmap": "roadmap",
    "timeline_roadmap": "roadmap",
    "comparison": "table_comparison",
    "cycle": "circular",
    "circular": "circular",
    "circular_cycle": "circular",
}


def _hex_to_rgb(hex_color: str):
    """Convert #RRGGBB to (R, G, B) tuple of 0-255 ints."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def _pptx_color(hex_color: str):
    """Return pptx RGBColor from hex string."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    r, g, b = _hex_to_rgb(hex_color)
    return RGBColor(r, g, b)


class SmartArtAgent(BaseAgent):
    """
    Draws SmartArt-style diagrams on PowerPoint slides.

    Input state keys:
        slide_contents (list): Per-slide content packages
        prs (Presentation): python-pptx Presentation object (passed by reference)
        brand_config (dict): Theme colors and fonts
        diagram_requests (list): [{slide_number, diagram_type, items, title}]

    Output (AgentResult.data):
        {
          diagram_slide_indices: {slide_number: pptx_slide_index},
          diagrams_created: int
        }
    """

    def __init__(self):
        super().__init__(
            name="SmartArtAgent",
            description="Creates SmartArt-style diagrams using python-pptx native shapes",
        )
        self._brand_colors = self._load_brand_colors()

    def _load_brand_colors(self) -> dict:
        config_path = Path(__file__).parent.parent / "config" / "brand_config.json"
        try:
            with open(config_path, "r") as f:
                cfg = json.load(f)
            return cfg.get("themes", {}).get("default", {})
        except Exception:
            return {
                "primary_color": "#002060",
                "secondary_color": "#00B0F0",
                "accent_color": "#FF6B35",
            }

    def _apply_theme_colors(self, theme: str):
        config_path = Path(__file__).parent.parent / "config" / "brand_config.json"
        try:
            with open(config_path, "r") as f:
                cfg = json.load(f)
            self._brand_colors = cfg.get("themes", {}).get(theme, self._brand_colors)
        except Exception:
            pass

    def draw_chevron_process(self, slide, items: list, title: str, colors: list):
        """
        Draw a horizontal chevron process flow.
        items: list of strings (process step labels, max 5 for readability)
        """
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        items = items[:5]  # max 5 chevrons
        n = len(items)
        if n == 0:
            return

        slide_width = Inches(10)
        margin_left = Inches(0.4)
        top = Inches(2.0)
        height = Inches(1.1)
        total_width = Inches(9.2)
        overlap = Inches(0.18)
        shape_width = (total_width + overlap * (n - 1)) / n

        for i, item in enumerate(items):
            color_hex = colors[i % len(colors)]
            left = margin_left + i * (shape_width - overlap)

            # Use pentagon/chevron shape
            from pptx.util import Inches
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            import pptx.oxml.ns as ns
            from pptx.oxml import parse_xml
            from lxml import etree

            # Add a rounded rectangle as chevron approximation
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
                left, top, shape_width, height
            )
            shape.fill.solid()
            r, g, b = _hex_to_rgb(color_hex)
            from pptx.dml.color import RGBColor
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
            shape.line.fill.background()

            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = item
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)

            # Step number badge
            badge = slide.shapes.add_shape(1, left + Inches(0.08), top - Inches(0.25), Inches(0.28), Inches(0.28))
            badge.fill.solid()
            badge.fill.fore_color.rgb = RGBColor(255, 255, 255)
            badge.line.color.rgb = RGBColor(r, g, b)
            badge.line.width = Pt(1.5)
            btf = badge.text_frame.paragraphs[0]
            btf.alignment = PP_ALIGN.CENTER
            brun = btf.add_run()
            brun.text = str(i + 1)
            brun.font.size = Pt(8)
            brun.font.bold = True
            brun.font.color.rgb = RGBColor(r, g, b)

    def draw_pyramid(self, slide, items: list, title: str, colors: list):
        """
        Draw a top-down pyramid diagram.
        items[0] = apex (most important), items[-1] = base.
        """
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        items = items[:5]
        n = len(items)
        if n == 0:
            return

        center_x = Inches(5.0)
        top_start = Inches(1.8)
        tier_height = Inches(0.85)
        max_width = Inches(6.0)

        for i, item in enumerate(items):
            tier_width = max_width * (i + 1) / n
            left = center_x - tier_width / 2
            top = top_start + i * tier_height
            color_hex = colors[i % len(colors)]
            r, g, b = _hex_to_rgb(color_hex)

            shape = slide.shapes.add_shape(1, left, top, tier_width, tier_height * 0.92)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
            shape.line.fill.background()

            tf = shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = item
            run.font.size = Pt(9.5)
            run.font.bold = (i == 0)
            run.font.color.rgb = RGBColor(255, 255, 255)

    def draw_circular_cycle(self, slide, items: list, title: str, colors: list):
        """
        Draw a circular cycle diagram with items as arc segments.
        items: list of step labels (3-6 recommended)
        """
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        import math

        items = items[:6]
        n = len(items)
        if n == 0:
            return

        cx = Inches(5.0)
        cy = Inches(3.8)
        radius = Inches(1.6)
        box_size = Inches(1.2)

        angle_step = 360 / n
        for i, item in enumerate(items):
            angle_deg = -90 + i * angle_step
            angle_rad = math.radians(angle_deg)
            left = cx + radius * math.cos(angle_rad) - box_size / 2
            top = cy + radius * math.sin(angle_rad) - box_size / 2

            color_hex = colors[i % len(colors)]
            r, g, b = _hex_to_rgb(color_hex)

            # Circle node
            shape = slide.shapes.add_shape(
                9,  # MSO_CONNECTOR_TYPE or oval — use oval (9)
                left, top, box_size, box_size
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
            shape.line.fill.background()

            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = item
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)

        # Center circle
        center_size = Inches(0.9)
        center_shape = slide.shapes.add_shape(
            9, cx - center_size / 2, cy - center_size / 2, center_size, center_size
        )
        center_shape.fill.solid()
        r, g, b = _hex_to_rgb(colors[0])
        center_shape.fill.fore_color.rgb = RGBColor(r, g, b)
        center_shape.line.fill.background()
        ctf = center_shape.text_frame.paragraphs[0]
        ctf.alignment = PP_ALIGN.CENTER
        crun = ctf.add_run()
        crun.text = "Cycle"
        crun.font.size = Pt(8)
        crun.font.bold = True
        crun.font.color.rgb = RGBColor(255, 255, 255)

    def draw_timeline_roadmap(self, slide, items: list, title: str, colors: list):
        """
        Draw a horizontal timeline roadmap with milestone markers.
        items: list of {label, date/period, description} dicts or strings
        """
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        items = items[:6]
        n = len(items)
        if n == 0:
            return

        # Timeline spine
        spine_top = Inches(3.2)
        spine_left = Inches(0.5)
        spine_width = Inches(9.0)
        spine_height = Inches(0.05)

        spine = slide.shapes.add_shape(1, spine_left, spine_top, spine_width, spine_height)
        spine.fill.solid()
        r, g, b = _hex_to_rgb(colors[0])
        spine.fill.fore_color.rgb = RGBColor(r, g, b)
        spine.line.fill.background()

        segment_width = spine_width / n
        for i, item in enumerate(items):
            label = item if isinstance(item, str) else item.get("label", f"Phase {i+1}")
            date_str = "" if isinstance(item, str) else item.get("date", item.get("period", ""))
            desc = "" if isinstance(item, str) else item.get("description", "")

            cx = spine_left + i * segment_width + segment_width / 2
            color_hex = colors[i % len(colors)]
            r, g, b = _hex_to_rgb(color_hex)

            # Milestone circle
            dot_size = Inches(0.22)
            dot = slide.shapes.add_shape(9, cx - dot_size / 2, spine_top - dot_size / 2, dot_size, dot_size)
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor(r, g, b)
            dot.line.fill.background()

            # Label box above
            box_h = Inches(0.55)
            box_w = Inches(1.5)
            box = slide.shapes.add_shape(1, cx - box_w / 2, spine_top - Inches(1.1), box_w, box_h)
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(r, g, b)
            box.line.fill.background()
            btf = box.text_frame
            btf.word_wrap = True
            bp = btf.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            brun = bp.add_run()
            brun.text = label
            brun.font.size = Pt(9)
            brun.font.bold = True
            brun.font.color.rgb = RGBColor(255, 255, 255)

            # Date below spine
            if date_str:
                dtb = slide.shapes.add_textbox(cx - Inches(0.75), spine_top + Inches(0.18), Inches(1.5), Inches(0.3))
                dtp = dtb.text_frame.paragraphs[0]
                dtp.alignment = PP_ALIGN.CENTER
                dtrun = dtp.add_run()
                dtrun.text = date_str
                dtrun.font.size = Pt(8)
                dtrun.font.color.rgb = RGBColor(120, 120, 120)
                dtrun.font.bold = False

    def map_content_to_diagram_type(self, content_type: str, visual_type: str) -> Optional[str]:
        """Determine diagram type from content_type or visual_type strings."""
        for key in (content_type, visual_type):
            if key and key.lower() in CONTENT_TYPE_MAP:
                return CONTENT_TYPE_MAP[key.lower()]
        return None

    def create_diagram_on_slide(
        self, slide, diagram_type: str, items: list, title: str, colors: list
    ):
        """Dispatch to the appropriate diagram drawer."""
        dispatchers = {
            "chevron": self.draw_chevron_process,
            "pyramid": self.draw_pyramid,
            "circular": self.draw_circular_cycle,
            "roadmap": self.draw_timeline_roadmap,
        }
        fn = dispatchers.get(diagram_type)
        if fn:
            try:
                fn(slide, items, title, colors)
                self.logger.info(f"Diagram '{diagram_type}' drawn on slide.")
            except Exception as e:
                self.logger.error(f"Diagram '{diagram_type}' failed: {e}", exc_info=True)
        else:
            self.logger.warning(f"Unknown diagram type: {diagram_type}")

    def execute(self, state: dict) -> AgentResult:
        """
        Process diagram requests and draw diagrams on the target slides.

        State keys consumed: diagram_requests, prs (Presentation object), intent_spec
        State keys produced: diagrams_created
        """
        try:
            diagram_requests = state.get("diagram_requests", [])
            if not diagram_requests:
                return AgentResult(
                    success=True,
                    data={"diagrams_created": 0},
                    metadata={"skipped": True, "reason": "no diagram requests"},
                )

            prs = state.get("prs")
            if prs is None:
                return AgentResult(
                    success=False,
                    data=None,
                    error="SmartArtAgent requires 'prs' (Presentation object) in state.",
                )

            intent_spec = state.get("intent_spec", {})
            theme = intent_spec.get("recommended_theme", "default")
            self._apply_theme_colors(theme)

            config_path = Path(__file__).parent.parent / "config" / "brand_config.json"
            try:
                with open(config_path, "r") as f:
                    cfg = json.load(f)
                colors = cfg.get("chart_colors", ["#002060", "#00B0F0", "#FF6B35", "#70AD47", "#FFC000"])
            except Exception:
                colors = ["#002060", "#00B0F0", "#FF6B35", "#70AD47", "#FFC000"]

            diagrams_created = 0
            slide_map = state.get("slide_index_map", {})

            for req in diagram_requests:
                slide_num = req.get("slide_number")
                diagram_type = req.get("diagram_type") or self.map_content_to_diagram_type(
                    req.get("content_type", ""), req.get("visual_type", "")
                )
                items = req.get("items", [])
                title = req.get("title", "")

                if not diagram_type or not items:
                    continue

                slide_idx = slide_map.get(slide_num)
                if slide_idx is None or slide_idx >= len(prs.slides):
                    self.logger.warning(f"Slide index for slide {slide_num} not found; skipping diagram.")
                    continue

                slide_obj = prs.slides[slide_idx]
                self.create_diagram_on_slide(slide_obj, diagram_type, items, title, colors)
                diagrams_created += 1

            self.logger.info(f"SmartArtAgent: {diagrams_created} diagrams created.")
            return AgentResult(
                success=True,
                data={"diagrams_created": diagrams_created},
                metadata={"diagrams_created": diagrams_created},
            )

        except Exception as e:
            self.logger.error(f"SmartArtAgent failed: {e}", exc_info=True)
            return AgentResult(success=False, data=None, error=str(e))
