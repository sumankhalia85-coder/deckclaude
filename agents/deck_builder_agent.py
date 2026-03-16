"""
DeckBuilderAgent — Assembles the final PowerPoint file using python-pptx.
Applies brand theme, creates slides from blueprint + content + visuals,
inserts charts (PNG), tables, and text with consistent formatting.
Font sizes: headline=28pt, body=16pt, caption=12pt. Margins: 0.5 inch.
"""

import json
import logging
import os
from pathlib import Path
from typing import Optional
from .base_agent import BaseAgent, AgentResult

logger = logging.getLogger(__name__)

# ---- Typography constants ----
FONT_HEADLINE = 28
FONT_BODY = 16
FONT_CAPTION = 12
FONT_SUBTITLE = 20
MARGIN = 0.5  # inches


def _hex_to_rgb(hex_color: str):
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


class DeckBuilderAgent(BaseAgent):
    """
    Builds the final .pptx file from all agent outputs.

    Input state keys:
        blueprint (dict): Slide structure and metadata
        slide_contents (list): Per-slide content packages
        visual_plan (list): Visual decisions per slide
        chart_files (dict): {slide_num: {path, chart_type, title}}
        image_files (dict): {slide_num: {path, attribution}}
        intent_spec (dict): Theme and metadata
        output_path (str, optional): Where to save the .pptx

    Output (AgentResult.data):
        {
          output_path: str,  # absolute path to .pptx file
          slide_count: int
        }
    """

    def __init__(self):
        super().__init__(
            name="DeckBuilderAgent",
            description="Builds final PowerPoint file using python-pptx",
        )
        self.output_dir = Path(os.getenv("OUTPUT_DIR", "./output"))
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def _load_theme(self, theme_name: str) -> dict:
        config_path = Path(__file__).parent.parent / "config" / "brand_config.json"
        try:
            with open(config_path, "r") as f:
                cfg = json.load(f)
            themes = cfg.get("themes", {})
            return themes.get(theme_name, themes.get("default", {}))
        except Exception:
            return {
                "primary_color": "#002060",
                "secondary_color": "#00B0F0",
                "accent_color": "#FF6B35",
                "background_color": "#FFFFFF",
                "text_color": "#1A1A2E",
                "font_heading": "Calibri",
                "font_body": "Calibri",
            }

    def _create_presentation(self):
        """Create a new pptx Presentation with widescreen 16:9 dimensions."""
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        return prs

    def _rgb_color(self, hex_color: str):
        from pptx.dml.color import RGBColor
        r, g, b = _hex_to_rgb(hex_color)
        return RGBColor(r, g, b)

    def _set_slide_background(self, slide, hex_color: str):
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb_color(hex_color)

    def _add_slide(self, prs, layout_index: int = 6):
        """Add a blank slide (layout index 6 = blank in most templates)."""
        try:
            layout = prs.slide_layouts[layout_index]
        except IndexError:
            layout = prs.slide_layouts[0]
        return prs.slides.add_slide(layout)

    def _add_text_box(
        self, slide, text: str, left, top, width, height,
        font_size: int, bold: bool = False, color_hex: str = "#1A1A2E",
        font_name: str = "Calibri", align_center: bool = False,
        wrap: bool = True,
    ):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        if align_center:
            p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = font_name
        r, g, b = _hex_to_rgb(color_hex)
        run.font.color.rgb = RGBColor(r, g, b)
        return txBox

    def _add_headline_bar(self, slide, headline: str, theme: dict, prs_width, prs_height):
        """Add a colored headline bar at the top of the slide."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        bar_height = Inches(1.1)
        bar = slide.shapes.add_shape(1, 0, 0, prs_width, bar_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb_color(theme.get("primary_color", "#002060"))
        bar.line.fill.background()

        tf = bar.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = headline
        run.font.size = Pt(FONT_HEADLINE)
        run.font.bold = True
        run.font.name = theme.get("font_heading", "Calibri")
        run.font.color.rgb = RGBColor(255, 255, 255)

        # Left margin inside bar
        from pptx.util import Inches
        tf.margin_left = Inches(MARGIN)
        tf.margin_top = Inches(0.15)

    def _add_bullet_content(
        self, slide, bullets: list, left, top, width, height,
        theme: dict, font_name: str, add_bullets: bool = True
    ):
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        if not bullets:
            return

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_top = Inches(0.05)

        for i, bullet in enumerate(bullets[:5]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(6)
            run = p.add_run()
            prefix = "• " if add_bullets else ""
            run.text = f"{prefix}{bullet}"
            run.font.size = Pt(FONT_BODY)
            run.font.name = font_name
            r, g, b = _hex_to_rgb(theme.get("text_color", "#1A1A2E"))
            run.font.color.rgb = RGBColor(r, g, b)

    def _add_chart_image(self, slide, chart_path: str, left, top, width, height):
        """Insert a chart PNG into the slide."""
        if not chart_path or not Path(chart_path).exists():
            return
        try:
            slide.shapes.add_picture(chart_path, left, top, width, height)
        except Exception as e:
            logger.warning(f"Failed to insert chart image {chart_path}: {e}")

    def _add_slide_number(self, slide, slide_num: int, total: int, prs_width, prs_height, theme: dict):
        """Add slide number in bottom-right corner."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        num_box = slide.shapes.add_textbox(
            prs_width - Inches(0.7), prs_height - Inches(0.35), Inches(0.6), Inches(0.3)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"{slide_num}"
        run.font.size = Pt(FONT_CAPTION)
        run.font.name = theme.get("font_body", "Calibri")
        run.font.color.rgb = self._rgb_color(theme.get("text_color", "#888888"))

    def _add_footer_line(self, slide, deck_title: str, prs_width, prs_height, theme: dict):
        """Add a thin footer bar at the bottom of the slide."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        footer_h = Inches(0.28)
        footer_top = prs_height - footer_h
        footer = slide.shapes.add_shape(1, 0, footer_top, prs_width, footer_h)
        footer.fill.solid()
        footer.fill.fore_color.rgb = self._rgb_color(theme.get("primary_color", "#002060"))
        footer.line.fill.background()

        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"  {deck_title}"
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(200, 200, 220)
        run.font.name = theme.get("font_body", "Calibri")

    def _build_title_slide(self, slide, slide_content: dict, image_info: Optional[dict],
                            theme: dict, prs_width, prs_height, deck_title: str):
        """Build a title slide with optional full-bleed background image."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        # Background: image or solid color
        if image_info and image_info.get("path") and Path(image_info["path"]).exists():
            try:
                slide.shapes.add_picture(image_info["path"], 0, 0, prs_width, prs_height)
            except Exception:
                self._set_slide_background(slide, theme.get("primary_color", "#002060"))
        else:
            self._set_slide_background(slide, theme.get("primary_color", "#002060"))

        # Dark overlay rectangle
        overlay = slide.shapes.add_shape(1, 0, 0, prs_width, prs_height)
        overlay.fill.solid()
        r, g, b = _hex_to_rgb(theme.get("primary_color", "#002060"))
        from pptx.dml.color import RGBColor as RC
        overlay.fill.fore_color.rgb = RC(r, g, b)
        from pptx.util import Pt as UPt
        # Set transparency via XML (python-pptx limitation — approximate)
        from lxml import etree
        solidFill = overlay.fill._xPr.solidFill
        alpha_elem = etree.SubElement(solidFill.find("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr") or solidFill, "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha")
        alpha_elem.set("val", "65000")  # ~65% opacity

        # Title text
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(10.5), Inches(2.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = slide_content.get("headline") or deck_title
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.name = theme.get("font_heading", "Calibri")
        run.font.color.rgb = RGBColor(255, 255, 255)

        # Subtitle / takeaway
        subtitle = slide_content.get("subtitle") or slide_content.get("takeaway", "")
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.6), Inches(9.0), Inches(0.7))
            stf = sub_box.text_frame
            sp = stf.paragraphs[0]
            sp.alignment = PP_ALIGN.LEFT
            srun = sp.add_run()
            srun.text = subtitle
            srun.font.size = Pt(FONT_SUBTITLE)
            srun.font.name = theme.get("font_body", "Calibri")
            srun.font.color.rgb = RGBColor(200, 220, 240)

        # Accent line
        line_shape = slide.shapes.add_shape(1, Inches(1.0), Inches(4.45), Inches(2.5), Inches(0.05))
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = self._rgb_color(theme.get("accent_color", "#FF6B35"))
        line_shape.line.fill.background()

    def _build_section_slide(self, slide, slide_content: dict, theme: dict, prs_width, prs_height):
        """Build a section divider slide."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        self._set_slide_background(slide, theme.get("secondary_color", "#00B0F0"))

        # Large section number
        sec_num = str(slide_content.get("slide_number", ""))
        if sec_num:
            num_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2.0), Inches(2.0))
            ntf = num_box.text_frame
            np_ = ntf.paragraphs[0]
            nrun = np_.add_run()
            nrun.text = sec_num
            nrun.font.size = Pt(96)
            nrun.font.bold = True
            nrun.font.color.rgb = RGBColor(255, 255, 255)
            nrun.font.name = theme.get("font_heading", "Calibri")

        # Section title
        title_box = slide.shapes.add_textbox(Inches(2.8), Inches(2.8), Inches(9.0), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = slide_content.get("headline", "Section")
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)

    def _build_thank_you_slide(self, slide, content: dict, theme: dict, prs_width, prs_height):
        """Build a 'Thank You' or Closing slide."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        self._set_slide_background(slide, theme.get("primary_color", "#002060"))

        # Large "Thank You" text
        title_box = slide.shapes.add_textbox(0, Inches(2.8), prs_width, Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = content.get("headline") or "Thank You"
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.name = theme.get("font_heading", "Calibri")

        # Contact text or concluding takeaway
        takeaway = content.get("takeaway", "")
        if takeaway:
            sub_box = slide.shapes.add_textbox(0, Inches(4.3), prs_width, Inches(1.0))
            stf = sub_box.text_frame
            sp = stf.paragraphs[0]
            sp.alignment = PP_ALIGN.CENTER
            srun = sp.add_run()
            srun.text = takeaway
            srun.font.size = Pt(18)
            srun.font.color.rgb = RGBColor(230, 230, 230)

    def _build_content_slide(
        self, slide, slide_content: dict, visual_plan_entry: dict,
        chart_path: Optional[str], image_info: Optional[dict],
        theme: dict, prs_width, prs_height, slide_num: int, total_slides: int, deck_title: str
    ):
        """Build a standard content slide with headline, bullets, and optional visual."""
        from pptx.util import Inches, Pt

        headline = slide_content.get("headline", "")
        bullets = slide_content.get("supporting_points", [])[:5]
        takeaway = slide_content.get("takeaway", "")
        data_callout = slide_content.get("data_callout")
        layout = visual_plan_entry.get("layout", "header_content")
        visual_type = visual_plan_entry.get("visual_type", "text_only")
        font_name = theme.get("font_body", "Calibri")

        # Background
        self._set_slide_background(slide, theme.get("background_color", "#FFFFFF"))

        # Headline bar
        self._add_headline_bar(slide, headline, theme, prs_width, prs_height)

        content_top = Inches(1.25)
        content_height = prs_height - Inches(1.85)
        left_margin = Inches(MARGIN)
        has_chart = bool(chart_path and Path(chart_path).exists())
        has_image = bool(image_info and image_info.get("path") and Path(image_info["path"]).exists())
        has_visual = has_chart or has_image

        if layout in ("two_column", "split_right", "split_left") or (visual_type in ("chart", "image") and has_visual) or (has_image and not has_chart):
            # Left: bullets, Right: visual
            col_width = prs_width * 0.44
            right_left = prs_width * 0.50
            right_width = prs_width * 0.46

            self._add_bullet_content(
                slide, bullets, left_margin, content_top, col_width - Inches(0.3), content_height, theme, font_name
            )

            # Data callout box on left if available
            if data_callout and data_callout.get("primary_number"):
                self._add_callout_box(
                    slide, data_callout, left_margin, prs_height - Inches(1.8),
                    col_width - Inches(0.3), Inches(0.7), theme
                )

            # Visual on right
            if has_chart:
                self._add_chart_image(slide, chart_path, right_left, content_top, right_width, content_height - Inches(0.1))
            elif has_image:
                try:
                    slide.shapes.add_picture(image_info["path"], right_left, content_top, right_width, content_height)
                except Exception as e:
                    logger.warning(f"Failed to insert image on slide {slide_num}: {e}")

        elif visual_type == "chart" and has_chart:
            # Chart-dominant: small bullets above, chart below
            bullet_height = Inches(0.9) if bullets else Inches(0)
            if bullets:
                self._add_bullet_content(
                    slide, bullets[:2], left_margin, content_top,
                    prs_width - Inches(1.0), bullet_height, theme, font_name
                )
            chart_top = content_top + bullet_height + Inches(0.1)
            chart_height = content_height - bullet_height - Inches(0.1)
            self._add_chart_image(slide, chart_path, left_margin, chart_top,
                                  prs_width - Inches(1.0), chart_height)

        else:
            # Text-only: full-width bullets
            self._add_bullet_content(
                slide, bullets, left_margin, content_top,
                prs_width - Inches(1.0), content_height - Inches(0.8), theme, font_name
            )
            # Takeaway in accent box at bottom
            if takeaway:
                self._add_takeaway_box(slide, takeaway, theme, prs_width, prs_height, font_name)

        # Footer and slide number
        self._add_footer_line(slide, deck_title, prs_width, prs_height, theme)
        self._add_slide_number(slide, slide_num, total_slides, prs_width, prs_height, theme)

    def _add_callout_box(self, slide, data_callout: dict, left, top, width, height, theme: dict):
        """Add a highlighted data callout box."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        box = slide.shapes.add_shape(1, left, top, width * 0.5, height)
        box.fill.solid()
        box.fill.fore_color.rgb = self._rgb_color(theme.get("accent_color", "#FF6B35"))
        box.line.fill.background()

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"{data_callout.get('primary_number', '')} {data_callout.get('unit', '')}"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.name = theme.get("font_heading", "Calibri")

    def _add_takeaway_box(self, slide, takeaway: str, theme: dict, prs_width, prs_height, font_name: str):
        """Add a subtle takeaway/insight box at the bottom of the slide."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        box_top = prs_height - Inches(1.2)
        box = slide.shapes.add_shape(1, Inches(MARGIN), box_top, prs_width - Inches(1.0), Inches(0.55))

        r, g, b = _hex_to_rgb(theme.get("primary_color", "#002060"))
        from pptx.dml.color import RGBColor as RC
        box.fill.solid()
        box.fill.fore_color.rgb = RC(min(r + 180, 255), min(g + 180, 255), min(b + 200, 255))
        box.line.color.rgb = self._rgb_color(theme.get("primary_color", "#002060"))
        box.line.width = Pt(0.5)

        tf = box.text_frame
        tf.margin_left = Inches(0.15)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Key takeaway: {takeaway}"
        run.font.size = Pt(FONT_CAPTION)
        run.font.italic = True
        run.font.name = font_name
        run.font.color.rgb = self._rgb_color(theme.get("primary_color", "#002060"))

    def build_presentation(self, state: dict) -> str:
        """
        Main build method: orchestrates slide creation for all slides.
        Returns the output file path.
        """
        from pptx.util import Inches

        blueprint = state.get("blueprint", {})
        slide_contents = state.get("slide_contents", [])
        visual_plan = state.get("visual_plan", [])
        chart_files = state.get("chart_files", {})
        image_files = state.get("image_files", {})
        intent_spec = state.get("intent_spec", {})
        output_path = state.get("output_path")

        theme_name = intent_spec.get("recommended_theme") or state.get("requested_theme") or "default"
        theme = self._load_theme(theme_name)
        self.logger.info(f"[DeckBuilder] Applying theme: '{theme_name}' → primary={theme.get('primary_color')}")
        deck_title = blueprint.get("deck_title") or intent_spec.get("deck_title", "Presentation")

        prs = self._create_presentation()
        prs_width = prs.slide_width
        prs_height = prs.slide_height

        # Build lookup maps
        content_by_num = {sc["slide_number"]: sc for sc in slide_contents}
        plan_by_num = {vp["slide_number"]: vp for vp in visual_plan}
        total_slides = len(blueprint.get("slides", []))

        slide_index_map = {}

        for bp_slide in blueprint.get("slides", []):
            snum = bp_slide["slide_number"]
            stype = bp_slide.get("slide_type", "content")
            content = content_by_num.get(snum, {})
            plan = plan_by_num.get(snum, {"visual_type": "text_only", "layout": "header_content"})
            chart_info = chart_files.get(snum)
            chart_path = chart_info.get("path") if chart_info else None
            image_info = image_files.get(snum)

            slide = self._add_slide(prs)
            slide_index_map[snum] = len(prs.slides) - 1

            if stype == "title":
                self._build_title_slide(slide, content, image_info, theme, prs_width, prs_height, deck_title)
            elif stype == "section":
                self._build_section_slide(slide, content, theme, prs_width, prs_height)
            elif stype == "thank_you":
                self._build_thank_you_slide(slide, content, theme, prs_width, prs_height)
            else:
                self._build_content_slide(
                    slide, content, plan, chart_path, image_info,
                    theme, prs_width, prs_height, snum, total_slides, deck_title
                )

            self.logger.info(f"Built slide {snum}/{total_slides}: [{stype}] {content.get('headline', '')[:60]}")

        # Determine output path
        if not output_path:
            safe_title = "".join(c if c.isalnum() or c in "-_ " else "_" for c in deck_title)[:40]
            output_path = str(self.output_dir / f"{safe_title.replace(' ', '_')}.pptx")

        prs.save(output_path)
        self.logger.info(f"Presentation saved: {output_path}")
        return output_path, slide_index_map

    def execute(self, state: dict) -> AgentResult:
        """
        Build the final PowerPoint presentation.

        State keys consumed: blueprint, slide_contents, visual_plan, chart_files,
                             image_files, intent_spec, output_path
        State keys produced: output_path, slide_index_map
        """
        try:
            blueprint = state.get("blueprint")
            slide_contents = state.get("slide_contents", [])

            if not blueprint:
                return AgentResult(
                    success=False,
                    data=None,
                    error="DeckBuilderAgent requires 'blueprint' in state.",
                )
            if not slide_contents:
                return AgentResult(
                    success=False,
                    data=None,
                    error="DeckBuilderAgent requires 'slide_contents' in state.",
                )

            output_path, slide_index_map = self.build_presentation(state)
            slide_count = len(state.get("blueprint", {}).get("slides", []))

            return AgentResult(
                success=True,
                data={
                    "output_path": output_path,
                    "slide_index_map": slide_index_map,
                    "slide_count": slide_count,
                },
                metadata={
                    "output_path": output_path,
                    "slide_count": slide_count,
                },
            )

        except Exception as e:
            self.logger.error(f"DeckBuilderAgent failed: {e}", exc_info=True)
            return AgentResult(success=False, data=None, error=str(e))
