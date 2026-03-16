"""
DeckClaude PowerPoint Utilities — python-pptx helper functions.

Provides: hex_to_rgb, set_shape_fill, add_text_to_shape, create_table_in_slide,
set_slide_background, add_slide_number, apply_brand_theme, add_accent_line,
set_paragraph_spacing, clear_slide_placeholders.
"""

import json
import logging
from pathlib import Path
from typing import List, Optional, Tuple, Union

logger = logging.getLogger(__name__)


def hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
    """
    Convert a hex color string to an (R, G, B) tuple of 0-255 integers.

    Args:
        hex_color: '#RRGGBB', 'RRGGBB', or '#RGB'

    Returns:
        (r, g, b) tuple

    Examples:
        hex_to_rgb('#002060') → (0, 32, 96)
        hex_to_rgb('FF6B35')  → (255, 107, 53)
    """
    hex_color = hex_color.lstrip("#").strip()
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    if len(hex_color) != 6:
        raise ValueError(f"Invalid hex color: '#{hex_color}'")
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return r, g, b


def rgb_color(hex_color: str):
    """Return a python-pptx RGBColor from a hex string."""
    from pptx.dml.color import RGBColor
    r, g, b = hex_to_rgb(hex_color)
    return RGBColor(r, g, b)


def set_shape_fill(shape, hex_color: str, transparency: float = 0.0):
    """
    Set a solid fill color on a pptx shape.

    Args:
        shape: python-pptx shape object
        hex_color: '#RRGGBB' color string
        transparency: 0.0 (opaque) to 1.0 (fully transparent)
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb_color(hex_color)

    if transparency > 0.0:
        # Encode as lumMod/lumOff or alpha in the XML
        alpha_val = int((1.0 - transparency) * 100000)
        try:
            solid_fill = shape.fill._xPr.solidFill
            from lxml import etree
            ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            clr_elem = solid_fill.find(f"{{{ns}}}srgbClr")
            if clr_elem is not None:
                alpha_elem = clr_elem.find(f"{{{ns}}}alpha")
                if alpha_elem is None:
                    alpha_elem = etree.SubElement(clr_elem, f"{{{ns}}}alpha")
                alpha_elem.set("val", str(alpha_val))
        except Exception as e:
            logger.debug(f"Could not set transparency: {e}")


def set_shape_border(shape, hex_color: str, width_pt: float = 1.0):
    """Set a solid border on a shape."""
    from pptx.util import Pt
    shape.line.color.rgb = rgb_color(hex_color)
    shape.line.width = Pt(width_pt)


def remove_shape_border(shape):
    """Remove (make invisible) a shape's border."""
    shape.line.fill.background()


def add_text_to_shape(
    shape,
    text: str,
    font_size: int = 12,
    bold: bool = False,
    italic: bool = False,
    color_hex: str = "#1A1A2E",
    font_name: str = "Calibri",
    align: str = "left",
    word_wrap: bool = True,
    space_before_pt: float = 0,
    space_after_pt: float = 0,
):
    """
    Add formatted text to a pptx shape's text frame (replaces existing content).

    Args:
        shape: python-pptx shape with text_frame
        text: Text content to add
        font_size: Font size in points
        bold/italic: Font style flags
        color_hex: Text color
        font_name: Font family name
        align: 'left' | 'center' | 'right' | 'justify'
        word_wrap: Enable text wrapping
        space_before_pt: Space before paragraph in points
        space_after_pt: Space after paragraph in points
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    pptx_align = align_map.get(align.lower(), PP_ALIGN.LEFT)

    tf = shape.text_frame
    tf.word_wrap = word_wrap

    # Clear existing content
    for para in tf.paragraphs[1:]:
        p_elem = para._p
        p_elem.getparent().remove(p_elem)

    p = tf.paragraphs[0]
    p.alignment = pptx_align
    if space_before_pt:
        p.space_before = Pt(space_before_pt)
    if space_after_pt:
        p.space_after = Pt(space_after_pt)

    # Clear existing runs
    for run in p.runs:
        r_elem = run._r
        r_elem.getparent().remove(r_elem)

    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    r, g, b = hex_to_rgb(color_hex)
    run.font.color.rgb = RGBColor(r, g, b)


def append_bullet_to_textframe(
    tf,
    text: str,
    font_size: int = 14,
    bold: bool = False,
    color_hex: str = "#1A1A2E",
    font_name: str = "Calibri",
    bullet_char: str = "•",
    indent_level: int = 0,
    space_before_pt: float = 4,
):
    """Append a new bullet paragraph to an existing text frame."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    p = tf.add_paragraph()
    p.space_before = Pt(space_before_pt)
    p.level = indent_level

    run = p.add_run()
    run.text = f"{bullet_char} {text}" if bullet_char else text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    r, g, b = hex_to_rgb(color_hex)
    run.font.color.rgb = RGBColor(r, g, b)


def create_table_in_slide(
    slide,
    headers: List[str],
    rows: List[List[str]],
    left,
    top,
    width,
    height,
    theme: dict = None,
    header_font_size: int = 11,
    body_font_size: int = 10,
    max_rows: int = 10,
):
    """
    Create a formatted data table on a slide.

    Args:
        slide: python-pptx slide object
        headers: Column header strings
        rows: List of row data (list of strings)
        left/top/width/height: Position and size (pptx Emu or Inches)
        theme: Brand config dict for colors
        header_font_size/body_font_size: Font sizes
        max_rows: Maximum rows to render
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if not headers or not rows:
        return

    theme = theme or {}
    primary_hex = theme.get("primary_color", "#002060")
    text_hex = theme.get("text_color", "#1A1A2E")
    font_name = theme.get("font_body", "Calibri")

    rows_to_show = rows[:max_rows]
    n_rows = len(rows_to_show) + 1  # +1 for header
    n_cols = len(headers)

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # Style header row
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        r, g, b = hex_to_rgb(primary_hex)
        cell.fill.fore_color.rgb = RGBColor(r, g, b)
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(header_text)
        run.font.size = Pt(header_font_size)
        run.font.bold = True
        run.font.name = font_name
        run.font.color.rgb = RGBColor(255, 255, 255)

    # Style data rows
    for row_idx, row_data in enumerate(rows_to_show, start=1):
        # Alternate row shading
        bg_hex = "#F2F5FA" if row_idx % 2 == 0 else "#FFFFFF"
        for col_idx in range(n_cols):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            r, g, b = hex_to_rgb(bg_hex)
            cell.fill.fore_color.rgb = RGBColor(r, g, b)

            cell_text = str(row_data[col_idx]) if col_idx < len(row_data) else ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = cell_text
            run.font.size = Pt(body_font_size)
            run.font.name = font_name
            r, g, b = hex_to_rgb(text_hex)
            run.font.color.rgb = RGBColor(r, g, b)

    return table


def set_slide_background(slide, hex_color: str):
    """Set a solid background color for a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color(hex_color)


def add_slide_number(
    slide,
    slide_num: int,
    prs_width,
    prs_height,
    font_name: str = "Calibri",
    color_hex: str = "#888888",
    font_size: int = 10,
):
    """Add a slide number indicator in the bottom-right corner."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    box = slide.shapes.add_textbox(
        prs_width - Inches(0.7),
        prs_height - Inches(0.35),
        Inches(0.6),
        Inches(0.28),
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(slide_num)
    run.font.size = Pt(font_size)
    run.font.name = font_name
    r, g, b = hex_to_rgb(color_hex)
    run.font.color.rgb = RGBColor(r, g, b)


def add_accent_line(
    slide,
    left,
    top,
    width,
    height=None,
    hex_color: str = "#FF6B35",
    vertical: bool = False,
):
    """Add a thin accent line (horizontal or vertical) to a slide."""
    from pptx.util import Inches, Pt

    if height is None:
        height = Inches(0.04) if not vertical else Inches(1.0)
    if vertical and width is None:
        width = Inches(0.04)

    line = slide.shapes.add_shape(1, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = rgb_color(hex_color)
    line.line.fill.background()


def apply_brand_theme(prs, theme: dict):
    """
    Apply brand theme settings to all slides in a presentation.
    Sets background color and ensures font consistency.

    Args:
        prs: python-pptx Presentation object
        theme: Dict from brand_config.json themes section
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    bg_hex = theme.get("background_color", "#FFFFFF")
    text_hex = theme.get("text_color", "#1A1A2E")
    font_name = theme.get("font_body", "Calibri")

    for slide in prs.slides:
        # Set background
        try:
            set_slide_background(slide, bg_hex)
        except Exception as e:
            logger.debug(f"Could not set slide background: {e}")

        # Apply font to all text frames
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.font.name:
                            run.font.name = font_name


def clear_slide_placeholders(slide):
    """Remove all default placeholder text boxes from a slide."""
    for shape in list(slide.placeholders):
        sp = shape._element
        sp.getparent().remove(sp)


def set_paragraph_spacing(paragraph, space_before_pt: float = 0, space_after_pt: float = 0, line_spacing: float = 1.0):
    """
    Set paragraph spacing for a pptx paragraph object.

    Args:
        paragraph: python-pptx paragraph
        space_before_pt: Space before paragraph in points
        space_after_pt: Space after paragraph in points
        line_spacing: Line spacing multiplier (1.0 = single, 1.5 = 1.5x)
    """
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    if space_before_pt:
        paragraph.space_before = Pt(space_before_pt)
    if space_after_pt:
        paragraph.space_after = Pt(space_after_pt)
    if line_spacing != 1.0:
        try:
            pPr = paragraph._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
            spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
            spcPct.set("val", str(int(line_spacing * 100000)))
        except Exception as e:
            logger.debug(f"Could not set line spacing: {e}")
